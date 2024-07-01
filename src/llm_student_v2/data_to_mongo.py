import os
import re
import PyPDF2
import pandas as pd
from openai import OpenAI
from pymongo import MongoClient
from pydub import AudioSegment
from pptx import Presentation

# splits audio files
def split_audio_file(file_path, minutes=8):
    """
    this function splits an audio file into 8 minute segments
    -Inputs:
    file_path: filepath to audio file
    minutes (int): specifies the chunk suze of the audio file. 12 minutes ~ 25mb
    -Output:
    segments (list): list of chunked audio files
    """
    six_minutes = minutes * 60 * 1000  # 6 minutes in milliseconds
    audio = AudioSegment.from_wav(file_path)
    segments = []
    
    for i in range(0, len(audio), six_minutes):
        # Extract 6 minute segments of the audio
        segments.append(audio[i:i + six_minutes])
    
    return segments

# checks if a wav file is over threshold mintues
def is_over_six_minutes(file_path, minutes=8):
    audio = AudioSegment.from_wav(file_path)
    return len(audio) > minutes * 60000  # duration in milliseconds

# transcribes audio files
def transcribe_audio(audio_segment):
    """
    This function takes an audio file and transcribes the text using openai's whisper model
    Inputs:
    -filepath: The path to the audio file (commonly supported files are wav, mp3, mp4 and a few more)
    """   
    client = OpenAI() # get openai api key from config
    # temp_dir = os.path.join(directory, "temp_dir")
    temp_dir = r"C:\Users\alipe\OneDrive\Desktop\llm_student_project\data\test\COMP500\temp_dir"
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir, exist_ok=True) # make temp_dir if it does not exist
    temp_filepath = os.path.join(temp_dir, 'temp_file.wav') # make a temp file directory
    audio_segment.export(temp_filepath, format="wav")  # Export segment to a temp file path
    with open(temp_filepath, "rb") as audio_file: # open the temp_filepath following open ai's documentation
        transcript = client.audio.transcriptions.create(model="whisper-1", # transcribe
                                                        file=audio_file,
                                                        response_format='text') 
    os.remove(temp_filepath)
    return transcript

# function to combine transcriptions
def combine_transcriptions(transcribed_texts):
    return ' '.join(transcribed_texts)

# processes audio_files
def process_audio_files(directory, files):
    """
    processes audio files found in a directory
    Inputs:
    -directory: path to a bunch of files
    -files (list): list of files found in that directory
    -openai_api_key: OpenAI api key required to acces whisper-1 model
    Outputs:
    -audio_documents (list): list of dictionaries structured for import into MongoDB

    """
    # A regex pattern for checking if a string ends with .wav
    file_pattern = r'.*\.(wav)$'
    audio_files = [f for f in files if re.search(file_pattern, f)] # take audio files from files
    documents = {}
    if len(audio_files) > 0:
        for wav_file in audio_files:
            wav_filepath = os.path.join(directory, wav_file)
            if is_over_six_minutes(wav_filepath):
                segments = split_audio_file(wav_filepath)
                print(f"Transcribing audio segments for {wav_file}...")
                transcriptions = [transcribe_audio(segment) for segment in segments] # , directory)
                print(f"Transcription for {wav_file} complete!\n")
                final_text = combine_transcriptions(transcriptions)
            else:
                # If file is not over 6 minutes, transcribe directly
                audio = AudioSegment.from_wav(wav_filepath)
                print(f"Short audio file, no split necessary. Transcribing audio file for {wav_file}...")
                final_text = transcribe_audio(audio)#, directory)
                print(f"Transcription for {wav_file} complete!\n")
            name, _ = os.path.splitext(wav_file)
            documents[name] = final_text
    else:
        print(f"No new audio files in this batch\n")
    classname = os.path.basename(directory) # this is the folder where the files reside
    audio_documents = format_for_mongodb(documents, classname)
    return audio_documents
    
# formats data for MongoDB insertion
def format_for_mongodb(lecture_text_dict, classname):
    """
    Format merged transcriptions into a structure suitable for MongoDB insertion.
    Args:
        merged_transcriptions (dict): Dictionary with titles as keys and concatenated texts as values.
        classname (str): The classname derived from the directory structure.
    Returns:
        list: A list of dictionaries formatted for MongoDB.
    """
    documents = []
    for title, text in lecture_text_dict.items():
        document = {
            "classname": classname,
            "title": title,
            "whole_document": {"text": text}
        }
        documents.append(document)
    return documents

# processes pptx files
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

    return " ".join(text_runs)

# processes txt and pdf files
def extract_text(directory, files):
    '''
    Extracts text from all pdf, txt files in a folder and appends individual file text into a dictionary where the name of the 
    file is the key and the text is the value. The driectory and file are separate because the update function in the main function below 
    feeds just the name of the new file.
    Inputs:
    -directory: The path to the folder containing the files.
    -files: this is name of the file in the path

    Outputs:
    -It outputs a dictionary containing the filename and the text
    '''
    print('extracting text from files...')
    full_paths = [os.path.join(directory, file) for file in files if file.endswith('.txt') or file.endswith('.pdf') or file.endswith('.pptx')]
    text_dict = {}
    for file in full_paths:
        lecture_title = os.path.splitext(os.path.basename(file))[0]
        print(f"lecture title: {lecture_title}, file: {file}")
        if file.lower().endswith('.pdf'):
            try:
                with open(file, 'rb') as pdf_file:
                    reader = PyPDF2.PdfReader(pdf_file)
                    text = "" # create empty string
                    for page in reader.pages:
                        page_text = page.extract_text() # extract text
                        if page_text:
                            text += page_text
                    text_dict[lecture_title] = text # exclude the last 4 characters which will the filename 
            except Exception as e:
                print(f"Error extracting text from this pdf, {file}: {e}")
        elif file.lower().endswith('.txt'):
            try:
                with open(file, 'rb') as txt_file:
                    text_dict[lecture_title] = txt_file.read() # load text into dict
            except Exception as e:  
                print(f"Error reading this txt, {file}: {e}")
        elif file.lower().endswith('.pptx'):
            try:
                text_dict[lecture_title] = extract_text_from_pptx(file) # extract text from pptx
            except Exception as e:
                print(f"Error reading this pptx file, {file}: {e}")
    print("Text extraction complete")
    #print(text_dict)
    return text_dict

# cleans extracted text from pdfs and text files       
def clean_text(text):
    if isinstance(text, bytes): # check if object is bytes
        text = text.decode("utf-8")  # convert text to from bytes object to str object
    
    # Regular expression pattern to match the unwanted text that looks similar to this:
    # 1/8/24, 10:32 AM about:blank\nabout:blank 1/1O
    pattern = r'\d{1,2}/\d{1,2}/\d{2,4},\s?\d{1,2}:\d{2}\s?[AaPp][Mm]\s?about:blank\s?'\
                        r'about:blank\s?\d/\d'
    # Removing the matched pattern
    cleaned_text = re.sub(pattern, '', text).strip()
    cleaned_text2 = re.sub(r'\n',' ', cleaned_text) # remove the \n
    return cleaned_text2

# processes pdf and txt files
def process_other_files(directory, files):
    """
    processes pdf and txt files and prepares formats them in json structure for mongo db import
    Input:
    -directory: The path to the folder containing the files.
    Output:
    -list: A list of dictionaries formatted for MongoDB.
    """
    lecture_text = extract_text(directory, files)
    text_clean = {name: clean_text(text) for name, text in lecture_text.items()}
    documents = format_for_mongodb(text_clean, os.path.basename(directory))
    return documents

# processes all files and saves them to a mongodb
def process_and_save_files(directory, files, database, mongo_uri):
    """
    processes all txt, pdf and wav files in a directory and imports into a mongo db
    Input:
    -directory: path to the files
    -files: new files found from find_new_files function
    -database: name of database, the collection remains the same
    -openai_api_key: openai api key
    -mongo_uri: mongo db connection string
    """
    print("processing audio files...")
    audio_files = process_audio_files(directory, files) # process audio files
    other_files = process_other_files(directory, files) # process txt and pdf files
    documents = audio_files + other_files # combine list of audio files and other files

    # connect to mongo db
    client = MongoClient(mongo_uri)
    db = client[database]
    collection = db['embeddings_v2']
    collection.insert_many(documents)
    print("Documents inserted into MongoDB successfully.")

# retrieves mongo docs from database
def get_mongo_docs(database, mongo_uri):
    """
    connects to a mongo database and returns all title fields in a pandas
    dataframe
    Inputs:
    -database (str): database name
    Outputs:
    -pandas Dataframe
    """
    try:
        client = MongoClient(mongo_uri)
        db = client[database]
        collection = db['embeddings_v2']
        documents = collection.find({}, {'title': 1, '_id': 0})
        return pd.DataFrame(list(documents))
    except Exception as e:
        print(f"Failed to connect to MongoDB with the following error:\n {e}")

# retrieves files from a directory
def get_files(directory, file_pattern):
    """
    gets all files from a directory that satisfy the file pattern filter
    Inputs:
    -directory: path to a folder containing new files you want to import
    -file_pattern (str): regex expression to be fed into re.search function
    Outputs:
    -pandas DataFrame containing a list of files in directory
    """
    files = [f for f in os.listdir(directory) if re.search(file_pattern, f)]
    if not files:
        return pd.DataFrame()
    return pd.DataFrame({'title': [f.split('.')[0] for f in files],
                         'extension': [f.split('.')[1] for f in files]})

# compare files form mongodb and directory
def compare_files(mongo_df, files_df):
    """
    Compares files to check if there are any new files in files_df that are 
    not in mongo_df.
    Inputs:
    -files_df: DataFrame containing files from folder
    -mongo_df: DataFrame containing file uploaded into MongoDB
    Outputs:
    -pandas DataFrame of files df
    """
    if files_df.empty:
        print(f"There are no .pdf, .txt, .wav files in folder")
        return None
    if not mongo_df.empty:
        files_df = files_df[~files_df['title'].isin(mongo_df['title'])].copy()
    files_df['combined'] = files_df['title'] + '.' + files_df['extension']
    return files_df

# finds new files in directory by comparing them to files in mongodb
def find_new_files(directory, database, mongo_uri):
    """
    This function searches a directory of files and find new files that haven't 
    been uploaded in the database
    Inputs:
    -directory: The directory path to the folder 
    -database: database to get already uploaded data (default is 'student_ai')
    Output:
    -python list containing new files in a directory
    """
    mongo_df = get_mongo_docs(database, mongo_uri)
    file_pattern = r'.*\.(pdf|txt|wav|pptx)$'
    files_df = get_files(directory, file_pattern)
    new_files_df = compare_files(mongo_df, files_df)
    if new_files_df is not None and not new_files_df.empty: # check if there is new files
        files_list = new_files_df['combined'].to_list()
        print(f"Found {len(new_files_df)} new file(s) to import in {directory}:")
        print(f"{files_list}\n")
        return files_list
    else:
        print("No new files to import")
        return []
    
# imports lectures
def import_lectures(directory, mongodb_name=None, mongo_uri=None):
    '''
    This function takes a file path to a folder containing a bunch of pdf, txt or wav files and loads them into a mongo db

    Inputs:
    -directory: file path to a folder containing pdf or txt files containing lecture transcripts. Name the folder
    containing the files after the class name for better organization in the database.
    -mongodb_name: the name of the database containing the collection 'embeddings'. The connection string is defined in config file
    '''
    try:
        new_files = find_new_files(directory, mongodb_name, mongo_uri) # find new files in a directory
    except FileNotFoundError as f:
        return f"Error {f} was returned please specify another directory"
    # cleans audio files in directory and then writes the cleaned files to the directory
    if len(new_files) > 0:
        print(f"directory: {directory}")
        process_and_save_files(directory, new_files, mongodb_name, mongo_uri) # processes and saves all new file to MongoDB
